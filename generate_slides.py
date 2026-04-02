#!/usr/bin/env python3
"""Generate Capco Profit Model & AI Impact PPTX slides."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors matching the website theme
BG_DARK = RGBColor(0x0A, 0x0A, 0x0F)
BG_CARD = RGBColor(0x16, 0x16, 0x1F)
BG_SURFACE = RGBColor(0x1E, 0x1E, 0x2A)
ACCENT = RGBColor(0x00, 0xBF, 0xA6)
ACCENT_LIGHT = RGBColor(0x00, 0xE6, 0xC8)
GOLD = RGBColor(0xF0, 0xC0, 0x40)
RED = RGBColor(0xFF, 0x6B, 0x6B)
WHITE = RGBColor(0xE8, 0xE8, 0xED)
MUTED = RGBColor(0x88, 0x88, 0xA0)
DIM = RGBColor(0x55, 0x55, 0x6A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=12, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def add_paragraph(tf, text, size=12, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, space_before=0, space_after=0, font_name="Calibri"):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def add_bullet(tf, text, size=11, color=MUTED, space_before=2, font_name="Calibri"):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(1)
    p.level = 0
    run = p.add_run()
    run.text = "  \u2022  " + text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = font_name
    return p


# ============================================================
# SLIDE 0: COVER
# ============================================================
slide0 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide0, BG_DARK)

# Accent line at top
line = slide0.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

# Badge
badge = add_text_box(slide0, Inches(1), Inches(1.5), Inches(5), Inches(0.5))
set_text(badge.text_frame, "PROMOTION CASE  \u2014  SUPPLEMENTAL SLIDES", size=11, bold=True, color=ACCENT)

# Title
title_box = add_text_box(slide0, Inches(1), Inches(2.3), Inches(10), Inches(2.2))
tf = title_box.text_frame
tf.word_wrap = True
set_text(tf, "Understanding Capco\u2019s", size=40, bold=True, color=WHITE, font_name="Calibri")
p2 = add_paragraph(tf, "Profit Model & AI-Driven Impact", size=40, bold=True, color=ACCENT, font_name="Calibri")

# Subtitle
sub_box = add_text_box(slide0, Inches(1), Inches(4.6), Inches(8), Inches(1))
tf = sub_box.text_frame
tf.word_wrap = True
set_text(tf, "How Capco turns engagements into profit \u2014 and how I leverage AI daily to deliver better client outcomes, accelerate delivery, and drive business development.", size=14, color=MUTED)

# Meta info
meta_box = add_text_box(slide0, Inches(1), Inches(5.9), Inches(10), Inches(0.8))
tf = meta_box.text_frame
set_text(tf, "Jessie Wang  \u00b7  Senior Consultant     |     Client: Jefferies \u2014 JefAI Program     |     Role: AI Product & Delivery Lead", size=12, color=DIM)

# Decorative accent shape (bottom right)
accent_rect = slide0.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.5), Inches(6.5), Inches(1.5), Inches(0.8))
accent_rect.fill.solid()
accent_rect.fill.fore_color.rgb = ACCENT
accent_rect.line.fill.background()
accent_rect.rotation = -10.0

# ============================================================
# SLIDE 1: CAPCO PROFIT MODEL
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide1, BG_DARK)

# Top accent line
line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

# Tag
tag = add_text_box(slide1, Inches(0.6), Inches(0.3), Inches(3), Inches(0.35))
set_text(tag.text_frame, "\u2014  BUSINESS ACUMEN", size=9, bold=True, color=ACCENT)

# Title
title = add_text_box(slide1, Inches(0.6), Inches(0.6), Inches(12), Inches(0.6))
tf = title.text_frame
tf.word_wrap = True
set_text(tf, "How Capco Makes Profit \u2014 and Where the Model Breaks Down", size=24, bold=True, color=WHITE, font_name="Calibri")

# --- 4 Flow Steps ---
flow_data = [
    ("01", "Opportunity Identification", "Client need surfaces via RFP or proactive pitch. Assess need, budget, authority, and timeline."),
    ("02", "Shaping & Pricing", "Build team pyramid (AO\u2192Partner). Plug into global pricing model: grade, location, bill rate, duration, holidays, PTO. Choose T&M or Fixed Price."),
    ("03", "Revenue Generation", "Contract signed with TCV. Fixed Price = even monthly payments. T&M = billable days \u00d7 rate/month. Revenue recognized as team delivers."),
    ("04", "Profit = The Arbitrage", "Margin = gap between client bill rate and Capco standard cost (salary + benefits). Target: 50%+ margin. Lower grades = higher margin."),
]

card_w = Inches(2.9)
card_h = Inches(1.55)
start_x = Inches(0.6)
gap = Inches(0.2)
y_flow = Inches(1.35)

for i, (num, heading, desc) in enumerate(flow_data):
    x = start_x + i * (card_w + gap)
    card = add_shape(slide1, x, y_flow, card_w, card_h, BG_CARD, RGBColor(0x2A, 0x2A, 0x3A))

    # Number badge
    num_box = add_text_box(slide1, x + Inches(0.15), y_flow + Inches(0.1), Inches(0.45), Inches(0.3))
    set_text(num_box.text_frame, num, size=9, bold=True, color=ACCENT)

    # Heading
    h_box = add_text_box(slide1, x + Inches(0.15), y_flow + Inches(0.35), card_w - Inches(0.3), Inches(0.3))
    set_text(h_box.text_frame, heading, size=11, bold=True, color=WHITE)

    # Description
    d_box = add_text_box(slide1, x + Inches(0.15), y_flow + Inches(0.65), card_w - Inches(0.3), Inches(0.85))
    d_box.text_frame.word_wrap = True
    set_text(d_box.text_frame, desc, size=9, color=MUTED)

# Arrow connectors between cards
for i in range(3):
    x = start_x + (i + 1) * (card_w + gap) - gap + Inches(0.02)
    arr = add_text_box(slide1, x, y_flow + Inches(0.6), gap, Inches(0.3))
    set_text(arr.text_frame, "\u25B6", size=12, color=ACCENT, alignment=PP_ALIGN.CENTER)

# --- Formula Bar ---
formula_y = Inches(3.05)
formula_bar = add_shape(slide1, Inches(0.6), formula_y, Inches(12.1), Inches(0.7), BG_CARD, RGBColor(0x00, 0x4D, 0x42))

formula_text = add_text_box(slide1, Inches(0.8), formula_y + Inches(0.08), Inches(11.7), Inches(0.55))
tf = formula_text.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER

parts = [
    ("Client Bill Rate ($X/day)", ACCENT_LIGHT, True),
    ("  \u2212  ", DIM, False),
    ("Capco Standard Cost ($Y/day)", GOLD, True),
    ("  =  ", DIM, False),
    ("Project Margin (Profit)", ACCENT, True),
    ("     |     ", DIM, False),
    ("Margin compresses at senior grades", RED, False),
]
for text, color, bold in parts:
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"

# --- 4 Breakdown Cards ---
breakdown_data = [
    ("\u26a0 Revenue Leakage", RED, [
        "Approved budget left un-billed (sick days, excess PTO)",
        "Client refuses milestone payments due to deliverable disputes",
        "T&M engagements under-tracked on timesheets",
    ]),
    ("\u26a0 Scope & Effort Overruns", RED, [
        "Delivery takes more hours than priced \u2014 bugs, complexity, rework",
        "Fixed price contracts absorb extra cost, eroding margin",
        "Insufficient buffer in original pricing model",
    ]),
    ("\u25b2 Market Pressure", GOLD, [
        "Competitors willing to price lower \u2014 forces rate drops",
        "Clients paying less for resource augmentation work",
        "Margins compressed over last 2\u20135 years industry-wide",
    ]),
    ("\u25b2 Utilization & Bench Cost", GOLD, [
        "Senior resources split across projects at part-time allocation",
        "Investment-rate staffing to protect relationships",
        "Bench time = cost with zero revenue offset",
    ]),
]

bd_w = Inches(5.95)
bd_h = Inches(1.35)
bd_y1 = Inches(3.95)
bd_y2 = Inches(5.4)
bd_positions = [
    (Inches(0.6), bd_y1), (Inches(6.7), bd_y1),
    (Inches(0.6), bd_y2), (Inches(6.7), bd_y2),
]

for idx, ((heading, accent_color, bullets), (bx, by)) in enumerate(zip(breakdown_data, bd_positions)):
    card = add_shape(slide1, bx, by, bd_w, bd_h, BG_CARD, RGBColor(0x2A, 0x2A, 0x3A))

    # Left accent border
    left_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx, by, Inches(0.05), bd_h)
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = accent_color
    left_bar.line.fill.background()

    # Content
    tb = add_text_box(slide1, bx + Inches(0.2), by + Inches(0.1), bd_w - Inches(0.35), bd_h - Inches(0.15))
    tb.text_frame.word_wrap = True
    set_text(tb.text_frame, heading, size=11, bold=True, color=WHITE)
    for bullet in bullets:
        add_bullet(tb.text_frame, bullet, size=9, color=MUTED, space_before=1)

# --- Jessie Insight Callout ---
callout_y = Inches(6.85)
callout = add_shape(slide1, Inches(0.6), callout_y, Inches(12.1), Inches(0.55), RGBColor(0x0D, 0x1F, 0x1C), ACCENT)

cb = add_text_box(slide1, Inches(0.8), callout_y + Inches(0.05), Inches(11.7), Inches(0.45))
cb.text_frame.word_wrap = True
tf = cb.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT

run = p.add_run()
run.text = "* My View as SC Leading AI at Jefferies:  "
run.font.size = Pt(10)
run.font.bold = True
run.font.color.rgb = ACCENT
run.font.name = "Calibri"

run2 = p.add_run()
run2.text = "I operate above my level daily \u2014 leading a 26-person team, owning the JefAI roadmap, presenting to C-level. Capco bills me at SC rate for PC/MP-level work. Billing at the appropriate grade = margin upside for Capco."
run2.font.size = Pt(10)
run2.font.color.rgb = MUTED
run2.font.name = "Calibri"

# ============================================================
# SLIDE 2: AI IN ACTION
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, BG_DARK)

# Top accent line
line = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

# Tag
tag = add_text_box(slide2, Inches(0.6), Inches(0.3), Inches(4), Inches(0.35))
set_text(tag.text_frame, "\u2014  AI-POWERED DELIVERY", size=9, bold=True, color=ACCENT)

# Title
title = add_text_box(slide2, Inches(0.6), Inches(0.6), Inches(12), Inches(0.6))
tf = title.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "How I Use AI to Deliver "
run.font.size = Pt(24)
run.font.bold = True
run.font.color.rgb = WHITE
run.font.name = "Calibri"
run2 = p.add_run()
run2.text = "Better Outcomes"
run2.font.size = Pt(24)
run2.font.bold = True
run2.font.color.rgb = ACCENT
run2.font.name = "Calibri"
run3 = p.add_run()
run3.text = " & Drive "
run3.font.size = Pt(24)
run3.font.bold = True
run3.font.color.rgb = WHITE
run3.font.name = "Calibri"
run4 = p.add_run()
run4.text = "Business Development"
run4.font.size = Pt(24)
run4.font.bold = True
run4.font.color.rgb = GOLD
run4.font.name = "Calibri"

# --- 6 AI Cards (3x2 grid) ---
ai_data = [
    ("\u26a1", "LOVABLE / FIGMA / CLAUDE", "Mockup to Interactive Prototype",
     "Rapidly transform concepts into clickable prototypes for stakeholder feedback before production code. Eliminates costly misalignment.",
     "\u2192 Weeks of dev time saved per feature cycle"),

    ("\u2318", "CLAUDE CODE / CODEX", "Full-Stack AI-Assisted Development",
     "Built the entire Jefferies AI Program Dashboard using Claude Code \u2014 a C-level reporting tool for project health, delivery tracking, and program metrics.",
     "\u2192 Information transparency for leadership & teams"),

    ("\u266b", "NOTEBOOKLM PODCAST", "Staying on the AI Frontier",
     "Use LLM-generated podcast briefings to stay current on AI developments \u2014 bringing cutting-edge insights to clients and aligning with JefAI roadmap.",
     "\u2192 Client always gets latest AI strategy guidance"),

    ("\u2b06", "INTERNAL CAPCO USE", "Codex & AI Tooling Internally",
     "Champion AI adoption within Capco \u2014 using Codex and Claude for internal projects, accelerating delivery velocity, and demonstrating ROI firm-wide.",
     "\u2192 Leading by example across the firm"),

    ("\u25cf", "CLAUDE CODE", "C-Level Reporting & Clarity",
     "AI-built dashboards bring information transparency to everyone \u2014 from ICs to CTO. Better project delivery tracking, clearer communication, faster decisions.",
     "\u2192 Clarity & alignment across all stakeholders"),

    ("\u25c6", "END-TO-END AI WORKFLOW", "The Full Picture",
     "From ideation to prototype to production \u2014 AI is embedded in every phase. A proven delivery model that saves time, reduces cost, and improves quality.",
     "\u2192 AI as a force multiplier, not a buzzword"),
]

ai_card_w = Inches(3.93)
ai_card_h = Inches(2.0)
ai_gap_x = Inches(0.2)
ai_gap_y = Inches(0.2)
ai_start_x = Inches(0.6)
ai_start_y = Inches(1.35)

for idx, (icon, tool_tag, heading, desc, outcome) in enumerate(ai_data):
    col = idx % 3
    row = idx // 3
    ax = ai_start_x + col * (ai_card_w + ai_gap_x)
    ay = ai_start_y + row * (ai_card_h + ai_gap_y)

    card = add_shape(slide2, ax, ay, ai_card_w, ai_card_h, BG_CARD, RGBColor(0x2A, 0x2A, 0x3A))

    # Icon
    icon_box = add_text_box(slide2, ax + Inches(0.15), ay + Inches(0.1), Inches(0.5), Inches(0.35))
    set_text(icon_box.text_frame, icon, size=18, color=WHITE)

    # Tool tag
    tag_box = add_text_box(slide2, ax + Inches(0.55), ay + Inches(0.12), ai_card_w - Inches(0.7), Inches(0.25))
    set_text(tag_box.text_frame, tool_tag, size=7, bold=True, color=ACCENT)

    # Heading
    h_box = add_text_box(slide2, ax + Inches(0.15), ay + Inches(0.45), ai_card_w - Inches(0.3), Inches(0.3))
    set_text(h_box.text_frame, heading, size=12, bold=True, color=WHITE)

    # Description
    d_box = add_text_box(slide2, ax + Inches(0.15), ay + Inches(0.75), ai_card_w - Inches(0.3), Inches(0.85))
    d_box.text_frame.word_wrap = True
    set_text(d_box.text_frame, desc, size=9, color=MUTED)

    # Outcome
    o_box = add_text_box(slide2, ax + Inches(0.15), ay + Inches(1.6), ai_card_w - Inches(0.3), Inches(0.3))
    set_text(o_box.text_frame, outcome, size=9, bold=True, color=ACCENT)

# --- Impact Banner ---
banner_y = Inches(5.8)
banner = add_shape(slide2, Inches(0.6), banner_y, Inches(12.1), Inches(1.1), BG_CARD, RGBColor(0x2A, 0x2A, 0x3A))

impact_stats = [
    ("10x", "Faster Prototyping"),
    ("26", "Engineering Team Led"),
    ("100%", "AI-Built Dashboard"),
    ("C-Suite", "Stakeholder Visibility"),
]

stat_w = Inches(3.0)
for i, (number, label) in enumerate(impact_stats):
    sx = Inches(0.6) + i * stat_w
    # Number
    n_box = add_text_box(slide2, sx, banner_y + Inches(0.15), stat_w, Inches(0.5))
    set_text(n_box.text_frame, number, size=26, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER, font_name="Calibri")
    # Label
    l_box = add_text_box(slide2, sx, banner_y + Inches(0.65), stat_w, Inches(0.3))
    set_text(l_box.text_frame, label, size=10, color=MUTED, alignment=PP_ALIGN.CENTER)

# --- Footer link ---
footer = add_text_box(slide2, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3))
set_text(footer.text_frame, "\u25c6  jessiewang0926.github.io/Jessie-Capco-Journey", size=9, color=DIM)

# ============================================================
# SAVE
# ============================================================
output_path = "/home/user/Jessie-Capco-Journey/Capco_Profit_Model_AI_Impact.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
